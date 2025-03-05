import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import matplotlib.pyplot as plt
import pandas as pd
from PIL import Image

def generate_presentation(title, slides):
    prs = Presentation()

    title_font = "Arial Black"
    content_font = "Calibri"
    theme_colors = {"title": "1F497D", "content": "4F81BD", "bullet": "FF5733"}

    slide_layout = prs.slide_layouts[0] 
    slide = prs.slides.add_slide(slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]

    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)

    subtitle_shape.text = "Auto-Generated Smart Presentation"
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(20)

    for slide_data in slides:
        if slide_data["type"] == "text":
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title_shape = slide.shapes.title
            content_shape = slide.shapes.placeholders[1]

            title_shape.text = slide_data["heading"]
            title_shape.text_frame.paragraphs[0].font.size = Pt(32)
            title_shape.text_frame.paragraphs[0].font.bold = True

            text_frame = content_shape.text_frame
            for i, point in enumerate(slide_data["points"]):
                p = text_frame.add_paragraph()
                p.text = f"• {point}"
                p.space_after = Pt(10)
                p.font.size = Pt(24)

        elif slide_data["type"] == "image":
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data["heading"]

            img_path = slide_data["image_path"]
            slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(6))

        elif slide_data["type"] == "chart":
            slide_layout = prs.slide_layouts[5]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data["heading"]

            fig, ax = plt.subplots()
            df = pd.DataFrame({"Category": ["A", "B", "C"], "Values": [20, 35, 50]})
            ax.bar(df["Category"], df["Values"], color=["red", "blue", "green"])
            ax.set_title("Sample Chart")

            chart_path = "chart.png"
            fig.savefig(chart_path)
            slide.shapes.add_picture(chart_path, Inches(1), Inches(1.5), width=Inches(6))

    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    pptx_bytes.seek(0)

    return pptx_bytes

st.set_page_config(page_title="📊 Smart Presentation Generator", layout="wide")
st.title("📊 Auto-Designed Smart Presentation Generator")

presentation_title = st.text_input("Enter Presentation Title", "My Presentation")

slide_type = st.selectbox("Select Slide Type", ["Text", "Image", "Chart"])

if "slides" not in st.session_state:
    st.session_state.slides = []

if slide_type == "Text":
    heading = st.text_input("Slide Heading")
    content = st.text_area("Enter bullet points (one per line)")
    if st.button("Add Slide"):
        st.session_state.slides.append({"type": "text", "heading": heading, "points": content.split("\n")})
        st.success("Text Slide Added!")

elif slide_type == "Image":
    heading = st.text_input("Slide Heading")
    image_file = st.file_uploader("Upload Image", type=["png", "jpg", "jpeg"])
    if image_file and st.button("Add Slide"):
        image_path = f"temp_{image_file.name}"
        with open(image_path, "wb") as f:
            f.write(image_file.getbuffer())
        st.session_state.slides.append({"type": "image", "heading": heading, "image_path": image_path})
        st.success("Image Slide Added!")

elif slide_type == "Chart":
    heading = st.text_input("Slide Heading")
    if st.button("Add Slide"):
        st.session_state.slides.append({"type": "chart", "heading": heading})
        st.success("Chart Slide Added!")

st.write("### Slides Preview:")
for i, slide in enumerate(st.session_state.slides):
    st.write(f"**Slide {i+1}: {slide['type'].capitalize()}** - {slide['heading']}")

    if slide["type"] == "text":
        st.write("\n".join(slide["points"]))

    elif slide["type"] == "image":
        img = Image.open(slide["image_path"])
        st.image(img, caption=slide["heading"], use_column_width=True)

    elif slide["type"] == "chart":
        df = pd.DataFrame({"Category": ["A", "B", "C"], "Values": [20, 35, 50]})
        fig, ax = plt.subplots()
        ax.bar(df["Category"], df["Values"], color=["red", "blue", "green"])
        ax.set_title("Sample Chart")
        st.pyplot(fig)

if st.button("Generate Presentation"):
    pptx_bytes = generate_presentation(presentation_title, st.session_state.slides)
    st.download_button("📥 Download Auto-Designed PPT", pptx_bytes, file_name="presentation.pptx")
 
